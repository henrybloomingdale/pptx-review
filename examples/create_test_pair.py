#!/usr/bin/env python3
"""
Create a pair of test .pptx files for testing pptx-review --diff.

Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def create_old():
    """Create test_old.pptx with 4 slides."""
    prs = Presentation()
    prs.core_properties.title = "Research Study"
    prs.core_properties.author = "Dr. Smith"

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Neuroimaging Study of Brain Connectivity"
    slide.placeholders[1].text = "Dr. Smith, Department of Neuroscience"

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Welcome everyone. Today I'll present our findings on brain connectivity."

    # Slide 2: Methods
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Methods"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Participants: 50 healthy adults (25M/25F)"
    p = tf.add_paragraph()
    p.text = "MRI Protocol: 3T Siemens scanner"
    p = tf.add_paragraph()
    p.text = "Analysis: FSL and FreeSurfer pipelines"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Methods slide - explain recruitment criteria"

    # Slide 3: Results
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Results"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Significant connectivity differences in DMN (p<0.01)"
    p = tf.add_paragraph()
    p.text = "Correlation with cognitive scores: r=0.45"
    p = tf.add_paragraph()
    p.text = "No significant age effects observed"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Key finding: DMN connectivity is altered"

    # Slide 4: Discussion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Discussion"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Findings consistent with prior literature"
    p = tf.add_paragraph()
    p.text = "Novel contribution: whole-brain analysis approach"
    p = tf.add_paragraph()
    p.text = "Clinical implications for early diagnosis"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Discuss limitations and future directions"

    prs.save("test_old.pptx")
    print("Created test_old.pptx (4 slides)")


def create_new():
    """Create test_new.pptx with modifications."""
    prs = Presentation()
    prs.core_properties.title = "Research Study - Revised"
    prs.core_properties.author = "Dr. Smith"

    # Slide 1: Title Slide (same structure, different notes)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Neuroimaging Study of Brain Connectivity"
    slide.placeholders[1].text = "Dr. Smith, Department of Neuroscience"

    # Changed notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Welcome. Today I'll present our UPDATED findings on brain connectivity patterns."

    # Slide 2: Methods (modified text)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Methods"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Participants: 75 healthy adults (38M/37F)"  # Changed from 50 to 75
    p = tf.add_paragraph()
    p.text = "MRI Protocol: 3T Siemens Prisma scanner"  # Added "Prisma"
    p = tf.add_paragraph()
    p.text = "Analysis: FSL, FreeSurfer, and CONN toolbox"  # Added CONN toolbox

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Methods slide - explain recruitment criteria"

    # Slide 3: Results (same as before)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Results"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Significant connectivity differences in DMN (p<0.01)"
    p = tf.add_paragraph()
    p.text = "Correlation with cognitive scores: r=0.45"
    p = tf.add_paragraph()
    p.text = "No significant age effects observed"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Key finding: DMN connectivity is altered"

    # Slide 4 (Discussion) is DELETED from new version

    # Slide 4 (was 5): Limitations - NEW slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Limitations"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Cross-sectional design limits causal inference"
    p = tf.add_paragraph()
    p.text = "Sample size may be insufficient for subgroup analyses"
    p = tf.add_paragraph()
    p.text = "Single-site study limits generalizability"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Be honest about limitations - reviewers will ask"

    prs.save("test_new.pptx")
    print("Created test_new.pptx (4 slides: modified slide 2, changed notes on slide 1, deleted Discussion, added Limitations)")


if __name__ == "__main__":
    create_old()
    create_new()
    print("\nDone! Now run:")
    print("  dotnet run --project ../PptxReview.csproj -- --diff test_old.pptx test_new.pptx")
    print("  dotnet run --project ../PptxReview.csproj -- --textconv test_new.pptx")
